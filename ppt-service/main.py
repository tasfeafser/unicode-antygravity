from fastapi import FastAPI, HTTPException
from fastapi.responses import JSONResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import base64
import json
from typing import Dict
import os

app = FastAPI()

class UnicodePPTGenerator:
    def __init__(self):
        self.template_cache = {}
    
    def replace_placeholders(self, prs: Presentation, content: Dict):
        """Replace placeholders in PPT with generated content"""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for key, value in content.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in shape.text:
                            shape.text = shape.text.replace(placeholder, str(value))
        
        return prs
    
    def add_code_block(self, slide, code: str, language: str):
        """Add syntax-highlighted code block to slide"""
        # Add code as text with monospace font
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = code
        text_frame.paragraphs[0].font.name = 'Courier New'
        text_frame.paragraphs[0].font.size = Pt(10)
        
        # Add language label
        label = slide.shapes.add_textbox(left, Inches(1.7), Inches(1), Inches(0.3))
        label.text_frame.text = f"📝 {language.upper()}"
    
    def add_diagram(self, slide, diagram_type: str, data: Dict):
        """Add simple ASCII diagram or reference to Mermaid"""
        # For Unicode platform, you can integrate with Mermaid.js
        # This is a placeholder for diagram generation
        pass

generator = UnicodePPTGenerator()

@app.post("/generate")
async def generate_ppt(request: Dict):
    try:
        # Decode template
        template_b64 = request.get('template')
        content = request.get('content', {})
        
        if not template_b64:
            raise HTTPException(400, "Template required")
        
        # Load template
        template_bytes = base64.b64decode(template_b64)
        prs = Presentation(BytesIO(template_bytes))
        
        # Replace placeholders
        prs = generator.replace_placeholders(prs, content)
        
        # Add code blocks if present
        if 'code_examples' in content:
            for code_example in content['code_examples']:
                # Add new slide for each code example
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                generator.add_code_block(slide, code_example['code'], code_example['language'])
        
        # Save to bytes
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        return JSONResponse({
            'pptx': base64.b64encode(output.getvalue()).decode('utf-8')
        })
    
    except Exception as e:
        raise HTTPException(500, str(e))

@app.get("/health")
async def health():
    return {"status": "healthy"}
